from pptx import Presentation
from pptx.util import Inches


class PresentationEngine:
    def create_outline(self, title: str, objective: str, audience: str, key_points: list[str]):
        return {
            "title": title,
            "objective": objective,
            "audience": audience,
            "slides": [
                {"slide": 1, "title": "Executive Summary", "content": [objective]},
                {"slide": 2, "title": "Key Insights", "content": key_points[:4]},
                {"slide": 3, "title": "Strategic Implications", "content": [
                    "Opportunities",
                    "Risks",
                    "Recommended next steps",
                ]},
            ],
            "style": "gamma-like executive storytelling structure",
        }

    def create_pptx(self, title: str, objective: str, audience: str, key_points: list[str], filename: str = "jarvis_presentation.pptx"):
        prs = Presentation()

        slide_layout = prs.slide_layouts[0]
        slide = prs.slides.add_slide(slide_layout)
        slide.shapes.title.text = title
        slide.placeholders[1].text = f"Audience: {audience}\nObjective: {objective}"

        bullet_layout = prs.slide_layouts[1]
        slide2 = prs.slides.add_slide(bullet_layout)
        slide2.shapes.title.text = "Key Insights"
        tf = slide2.placeholders[1].text_frame
        for idx, point in enumerate(key_points[:6]):
            if idx == 0:
                tf.text = point
            else:
                p = tf.add_paragraph()
                p.text = point

        slide3 = prs.slides.add_slide(bullet_layout)
        slide3.shapes.title.text = "Recommended Next Steps"
        tf3 = slide3.placeholders[1].text_frame
        tf3.text = "Validate assumptions with real data."
        p = tf3.add_paragraph()
        p.text = "Protect downside before scaling exposure."
        p = tf3.add_paragraph()
        p.text = "Align strategy with wealth, family, health, and time."

        output_path = f"generated_reports/{filename}"
        prs.save(output_path)
        return {
            "status": "created",
            "file": output_path,
        }
